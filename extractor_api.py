import io
import logging
import os
import subprocess
import tempfile
from pathlib import Path
from typing import List, Optional
import re
import asyncio

import httpx
from fastapi import FastAPI, HTTPException, Response, Body
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, HttpUrl
from weasyprint import HTML
from pptx import Presentation

from graph_utils import (
    download_file_from_graph,
    upload_file_to_graph,
    list_folder_children,
    get_item_name,
    startup_graph_client,
    close_graph_client,
)


app = FastAPI(title="PowerPoint Notes Extraction API")

# Enable CORS for all origins (use more restrictive settings in production)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# Health check endpoint for Azure App Service
@app.get("/health")
def health() -> dict[str, str]:
    """Simple health check returning application status."""
    return {"status": "ok"}


logger = logging.getLogger("pptx_extractor")
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s [%(name)s] %(message)s",
)

# Allow the request timeout to be configured via an environment variable.
# Defaults to 60 seconds if not provided.
TIMEOUT = float(os.environ.get("REQUEST_TIMEOUT", "60"))
CONNECT_TIMEOUT = float(os.environ.get("CONNECT_TIMEOUT", "10"))
DOWNLOAD_CONCURRENCY = int(os.environ.get("DOWNLOAD_CONCURRENCY", "5"))

# httpx requires all timeout parameters be specified when using custom values.
# Create a reusable configuration object shared by the client and per-request
# calls so that connection and read timeouts are explicit.
# Use a single timeout value for all operations to avoid misconfiguration
HTTPX_TIMEOUT = httpx.Timeout(TIMEOUT)

# Reusable HTTP client for outbound requests - created at startup
http_client: httpx.AsyncClient | None = None


@app.on_event("startup")
async def startup_event() -> None:
    """Create shared HTTP clients."""
    global http_client
    limits = httpx.Limits(max_connections=100)
    http_client = httpx.AsyncClient(timeout=HTTPX_TIMEOUT, limits=limits)
    await startup_graph_client()


@app.on_event("shutdown")
async def shutdown_event() -> None:
    """Close shared HTTP clients."""
    if http_client is not None:
        await http_client.aclose()
    await close_graph_client()

# External tool locations can be overridden via environment variables
FFPROBE_BIN = os.environ.get("FFPROBE_BIN", "ffprobe")
LIBREOFFICE_BIN = os.environ.get("LIBREOFFICE_BIN", "libreoffice")


class ExtractRequest(BaseModel):
    file_url: HttpUrl
    file_name: str


class SlideData(BaseModel):
    slide_number: int
    title_text: Optional[str] = None
    notes_text: Optional[str] = None


class ExtractResponse(BaseModel):
    filename: str
    slide_count: int
    slides: List[SlideData]


class CombineRequest(BaseModel):
    drive_id: str
    folder_id: str
    pptx_file_id: str


class CombineResponse(BaseModel):
    status: str
    video_filename: str
    upload_url: str



# HTML to PDF conversion now accepts raw HTML bytes rather than text.



def _extract_slides(presentation: Presentation) -> List[SlideData]:
    """Return slide metadata from a Presentation."""
    slides = []
    for idx, slide in enumerate(presentation.slides, start=1):
        title = slide.shapes.title.text if slide.shapes.title else None
        notes = (
            slide.notes_slide.notes_text_frame.text
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame
            else None
        )
        slides.append(SlideData(slide_number=idx, title_text=title, notes_text=notes))
    return slides


def _parse_slide_number(path: Path) -> int:
    """Return the numeric slide number embedded in a filename."""
    match = re.search(r"(\d+)", path.stem)
    if match:
        return int(match.group(1))
    logger.warning("Unexpected slide filename: %s", path.name)
    return 0


def _html_to_pdf_bytes(html: bytes) -> bytes:
    """Return PDF data generated from HTML bytes."""
    buf = io.BytesIO()
    try:
        HTML(string=html.decode()).write_pdf(target=buf)
    except Exception as exc:  # pylint: disable=broad-except
        raise ValueError("pdf generation failed") from exc
    return buf.getvalue()


def get_audio_duration(path: Path) -> float:
    """Return audio duration in seconds using ffprobe."""
    cmd = [
        FFPROBE_BIN,
        "-v",
        "error",
        "-show_entries",
        "format=duration",
        "-of",
        "default=noprint_wrappers=1:nokey=1",
        str(path),
    ]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, check=True)
    except FileNotFoundError as exc:
        logger.exception("ffprobe not found")
        raise HTTPException(status_code=500, detail="ffprobe is not installed") from exc
    except subprocess.CalledProcessError as exc:
        logger.exception("ffprobe failed")
        raise HTTPException(
            status_code=500,
            detail="Audio metadata extraction failed",
        ) from exc
    return float(result.stdout.strip())


async def calculate_slide_durations(audio_paths: List[Path]) -> List[float]:
    """Return per-slide durations including crossfade padding."""
    tasks = [asyncio.to_thread(get_audio_duration, p) for p in audio_paths]
    mp3_durations = await asyncio.gather(*tasks)
    return [duration + 2.0 for duration in mp3_durations]


async def run_cmd(cmd: List[str]) -> None:
    """Run an external command asynchronously and raise on failure."""
    proc = await asyncio.create_subprocess_exec(
        *cmd,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
    )
    stdout, stderr = await proc.communicate()
    if proc.returncode != 0:
        raise subprocess.CalledProcessError(
            proc.returncode, cmd, output=stdout, stderr=stderr
        )


@app.post("/extract", response_model=ExtractResponse)
async def extract_notes(request: ExtractRequest):
    url = str(request.file_url)
    logger.info("Extraction requested for %s", url)

    try:
        response = await http_client.get(
            url, timeout=HTTPX_TIMEOUT, follow_redirects=True
        )
        response.raise_for_status()
        logger.debug("Downloaded %d bytes", len(response.content))
    except (httpx.RequestError, httpx.HTTPStatusError) as exc:
        logger.exception("Failed to download file from %s", url)
        raise HTTPException(
            status_code=400, detail=f"Unable to download file: {exc}"
        ) from exc

    content_type = response.headers.get("Content-Type", "")
    logger.debug("Content-Type received: %s", content_type)
    if (
        content_type
        and "presentation" not in content_type
        and "ppt" not in content_type
    ):
        logger.warning("Unsupported content type: %s", content_type)
        raise HTTPException(status_code=422, detail="Only .pptx files are supported")

    pptx_bytes = io.BytesIO(response.content)
    try:
        presentation = Presentation(pptx_bytes)
    except Exception as exc:  # pylint: disable=broad-except
        logger.exception("Failed to parse PowerPoint file")
        raise HTTPException(status_code=422, detail="Invalid .pptx file") from exc

    slides_data = _extract_slides(presentation)

    logger.info("Successfully extracted %d slides", len(slides_data))

    filename = request.file_name
    return ExtractResponse(
        filename=filename,
        slide_count=len(slides_data),
        slides=slides_data,
    )


@app.post("/html-to-pdf/async")
async def html_to_pdf_async(html: bytes = Body(...)) -> Response:
    """Generate a PDF from provided HTML bytes asynchronously."""
    try:
        pdf_bytes = await asyncio.to_thread(_html_to_pdf_bytes, html)
    except Exception as exc:  # pylint: disable=broad-except
        logger.exception("PDF generation failed")
        raise HTTPException(status_code=500, detail="PDF generation failed") from exc
    return Response(content=pdf_bytes, media_type="application/pdf")


@app.post("/html-to-pdf")
def html_to_pdf(html: bytes = Body(...)) -> Response:
    """Generate a PDF from provided HTML bytes synchronously."""
    try:
        pdf_bytes = _html_to_pdf_bytes(html)
    except Exception as exc:  # pylint: disable=broad-except
        logger.exception("PDF generation failed")
        raise HTTPException(status_code=500, detail="PDF generation failed") from exc
    return Response(content=pdf_bytes, media_type="application/pdf")


@app.post("/combine", response_model=CombineResponse)
async def combine_presentation(request: CombineRequest):
    """Combine a PPTX presentation with per-slide audio into a video."""
    logger.info(
        "Combining PPTX %s in drive %s with audio from folder %s",
        request.pptx_file_id,
        request.drive_id,
        request.folder_id,
    )

    drive_id = request.drive_id
    folder_id = request.folder_id
    pptx_id = request.pptx_file_id

    try:
        pptx_bytes = await download_file_from_graph(drive_id, pptx_id)
        pptx_name = await get_item_name(drive_id, pptx_id)
    except Exception as exc:  # pylint: disable=broad-except
        status = getattr(getattr(exc, "response", None), "status_code", None)
        if status is not None:
            logger.error("Failed to download PPTX from Graph: HTTP %s", status)
        logger.exception("Failed to download PPTX from Graph")
        raise HTTPException(
            status_code=400,
            detail=f"Unable to download PPTX: {exc}",
        ) from exc

    with tempfile.TemporaryDirectory() as tmpdir:
        tmp_path = Path(tmpdir)
        pptx_path = tmp_path / "presentation.pptx"
        pptx_path.write_bytes(pptx_bytes)

        # Retrieve MP3 metadata from the folder
        try:
            children = await list_folder_children(drive_id, folder_id)
        except Exception as exc:  # pylint: disable=broad-except
            logger.exception("Failed to list folder contents")
            raise HTTPException(
                status_code=400, detail="Unable to list folder"
            ) from exc

        audio_items = [
            item for item in children if item.get("name", "").lower().endswith(".mp3")
        ]
        if not audio_items:
            raise HTTPException(status_code=400, detail="No MP3 files found")

        # Sort by slide number based on filename pattern slide_<n>.mp3
        audio_items.sort(key=lambda x: int(Path(x["name"]).stem.split("_")[-1]))

        audio_paths: List[Path] = []

        semaphore = asyncio.Semaphore(DOWNLOAD_CONCURRENCY)

        async def fetch_audio(item: dict) -> Path:
            async with semaphore:
                audio_bytes = await asyncio.wait_for(
                    download_file_from_graph(drive_id, item["id"]), timeout=TIMEOUT
                )
                audio_path = tmp_path / item["name"]
                audio_path.write_bytes(audio_bytes)
                await asyncio.to_thread(get_audio_duration, audio_path)
                return audio_path

        results = await asyncio.gather(*(fetch_audio(it) for it in audio_items))
        for path in results:
            audio_paths.append(path)

        slides_dir = tmp_path / "slides"
        slides_dir.mkdir(exist_ok=True)
        try:
            await run_cmd(
                [
                    LIBREOFFICE_BIN,
                    "--headless",
                    "--convert-to",
                    "png",
                    str(pptx_path),
                    "--outdir",
                    str(slides_dir),
                ]
            )
        except FileNotFoundError as exc:
            logger.exception("libreoffice not found")
            raise HTTPException(
                status_code=500,
                detail="libreoffice is not installed",
            ) from exc
        except subprocess.CalledProcessError as exc:
            logger.exception(
                "Slide image conversion failed: %s",
                exc.stderr.decode(errors="replace") if exc.stderr else ""
            )
            raise HTTPException(
                status_code=500,
                detail="PPTX conversion failed",
            ) from exc

        image_files = sorted(
            slides_dir.glob("*.png"),
            key=_parse_slide_number,
        )

        if len(image_files) != len(audio_paths):
            logger.warning(
                "Slide count mismatch: %d images vs %d audio files",
                len(image_files),
                len(audio_paths),
            )
            raise HTTPException(status_code=400, detail="Slide count mismatch")

        output_path = tmp_path / f"{Path(pptx_name).stem}.mp4"
        cmd = [
            "ffmpeg",
            "-y",
            str(output_path),
        ]

        try:
            await run_cmd(cmd)
        except FileNotFoundError as exc:
            logger.exception("ffmpeg not found")
            raise HTTPException(
                status_code=500,
                detail="ffmpeg is not installed",
            ) from exc
        except subprocess.CalledProcessError as exc:
            logger.exception("ffmpeg failed")
            raise HTTPException(
                status_code=500, detail="Video generation failed"
            ) from exc

        try:
            video_bytes = output_path.read_bytes()
            upload_url = await upload_file_to_graph(
                drive_id, folder_id, output_path.name, video_bytes
            )
        except Exception as exc:  # pylint: disable=broad-except
            logger.exception("Upload of generated video failed")
            raise HTTPException(status_code=500, detail="Upload failed") from exc

    return CombineResponse(
        status="success",
        video_filename=output_path.name,
        upload_url=upload_url,
    )


# For local development
if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000)
